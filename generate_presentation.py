#!/usr/bin/env python3
"""
Generate a Google Slides-style presentation from the AI Infrastructure Onboarding README files.
Uses GCP brand colors, styling, and iconography.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── GCP Brand Colors ──
GCP_BLUE = RGBColor(0x42, 0x85, 0xF4)       # Google Blue
GCP_RED = RGBColor(0xEA, 0x43, 0x35)         # Google Red
GCP_YELLOW = RGBColor(0xFB, 0xBC, 0x05)      # Google Yellow
GCP_GREEN = RGBColor(0x34, 0xA8, 0x53)       # Google Green
GCP_DARK_BLUE = RGBColor(0x17, 0x4E, 0xA6)   # GCP Dark Blue
GCP_PRIMARY = RGBColor(0x1A, 0x73, 0xE8)     # GCP Primary Blue
GCP_DARK = RGBColor(0x20, 0x21, 0x24)        # Google Dark Gray
GCP_MEDIUM_GRAY = RGBColor(0x5F, 0x63, 0x68) # Medium Gray
GCP_LIGHT_GRAY = RGBColor(0xF1, 0xF3, 0xF4)  # Light Gray Background
GCP_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GCP_ACCENT_TEAL = RGBColor(0x00, 0xBF, 0xA5) # Teal accent
GCP_SECTION_BG = RGBColor(0x17, 0x4E, 0xA6)  # Section divider background

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_gcp_bar(slide, prs):
    """Add the iconic 4-color Google bar at the bottom of a slide."""
    bar_height = Inches(0.06)
    bar_y = SLIDE_HEIGHT - bar_height
    colors = [GCP_BLUE, GCP_RED, GCP_YELLOW, GCP_GREEN]
    bar_width = SLIDE_WIDTH / 4
    for i, color in enumerate(colors):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(bar_width * i)),
            bar_y,
            bar_width,
            bar_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def add_header_bar(slide, color=None):
    """Add a colored header bar at the top."""
    if color is None:
        color = GCP_PRIMARY
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_gcp_icon(slide, icon_text, x, y, size=Inches(0.6), color=GCP_PRIMARY):
    """Add a circular icon with text (simulating GCP product icons)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x, y, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(int(size / Inches(1) * 14))
    p.font.color.rgb = GCP_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def add_hexagon_icon(slide, x, y, size=Inches(0.55), color=GCP_PRIMARY):
    """Add a hexagonal shape (GCP-style)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        x, y, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_slide_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=GCP_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=GCP_DARK, bullet_color=GCP_PRIMARY, spacing=Pt(6)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, title, body, icon_text="",
             accent_color=GCP_PRIMARY):
    """Add a material-design-style card with shadow effect."""
    # Card background (shadow)
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        left + Inches(0.04), top + Inches(0.04), width, height
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0xDA, 0xDC, 0xE0)
    shadow.line.fill.background()

    # Card main
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = GCP_WHITE
    card.line.color.rgb = RGBColor(0xDA, 0xDC, 0xE0)
    card.line.width = Pt(1)

    # Accent bar at top of card
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.05), top + Inches(0.05),
        width - Inches(0.1), Inches(0.05)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2),
                 width - Inches(0.4), Inches(0.4),
                 title, font_size=14, color=GCP_DARK, bold=True)

    # Body
    add_text_box(slide, left + Inches(0.2), top + Inches(0.55),
                 width - Inches(0.4), height - Inches(0.7),
                 body, font_size=11, color=GCP_MEDIUM_GRAY)


def create_title_slide(prs):
    """Slide 1: Title slide with GCP branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, GCP_DARK_BLUE)

    # Large decorative shapes
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.4), SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GCP_PRIMARY
    shape.line.fill.background()

    # Decorative circles
    for (x, y, s), col in [
        ((Inches(10.5), Inches(0.5), Inches(1.2)), GCP_RED),
        ((Inches(11.2), Inches(1.5), Inches(0.8)), GCP_YELLOW),
        ((Inches(10.8), Inches(2.5), Inches(0.5)), GCP_GREEN),
        ((Inches(11.5), Inches(0.3), Inches(0.4)), GCP_BLUE),
    ]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, s, s)
        circle.fill.solid()
        circle.fill.fore_color.rgb = col
        circle.line.fill.background()
        circle.fill.fore_color.brightness = 0.15

    # Google Cloud logo placeholder
    add_text_box(slide, Inches(1.2), Inches(1.0), Inches(5), Inches(0.6),
                 "☁  Google Cloud", font_size=22, color=GCP_WHITE, bold=True)

    # Title
    add_text_box(slide, Inches(1.2), Inches(2.0), Inches(9), Inches(1.5),
                 "AI Infrastructure Onboarding", font_size=44, color=GCP_WHITE, bold=True)

    # Subtitle
    add_text_box(slide, Inches(1.2), Inches(3.6), Inches(9), Inches(1.0),
                 "A comprehensive guide to Google Cloud AI infrastructure capabilities\nfor Public Sector customers",
                 font_size=20, color=RGBColor(0xAE, 0xCB, 0xFA))

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2), Inches(5.0), Inches(3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GCP_YELLOW
    line.line.fill.background()

    # Description
    add_text_box(slide, Inches(1.2), Inches(5.3), Inches(8), Inches(1.2),
                 "From quota requests and reservations, through deploying workloads on\n"
                 "GKE AI Hypercompute, to monitoring TPU health and performance.",
                 font_size=14, color=RGBColor(0xBD, 0xC1, 0xC6))

    # Bottom bar
    colors = [GCP_BLUE, GCP_RED, GCP_YELLOW, GCP_GREEN]
    bar_h = Inches(0.06)
    for i, c in enumerate(colors):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Emu(int(SLIDE_WIDTH / 4 * i)),
                                   SLIDE_HEIGHT - bar_h, SLIDE_WIDTH // 4, bar_h)
        s.fill.solid()
        s.fill.fore_color.rgb = c
        s.line.fill.background()


def create_agenda_slide(prs):
    """Slide 2: Agenda / Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, GCP_WHITE)
    add_header_bar(slide)
    add_gcp_bar(slide, prs)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "Onboarding Journey", font_size=32, color=GCP_DARK, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
                 "Five phases from planning to production", font_size=16, color=GCP_MEDIUM_GRAY)

    phases = [
        ("01", "Foundational Tools", "Accelerator guide, agentic coders,\nquota & reservations", GCP_BLUE),
        ("02", "Core Infrastructure", "Networking, disk images, storage,\ndata pipelines", GCP_RED),
        ("03", "Deploying Workloads", "DWS, GKE, Vertex AI, Cluster\nToolkit, XPK, Slurm", GCP_YELLOW),
        ("04", "Monitoring", "TPU observability, XProf,\nalerts & dashboards", GCP_GREEN),
        ("05", "Billing & Usage", "Budgets, quotas, CUDs,\ncost optimization", GCP_PRIMARY),
    ]

    start_x = Inches(0.5)
    card_width = Inches(2.35)
    card_height = Inches(3.8)
    gap = Inches(0.15)

    for i, (num, title, desc, color) in enumerate(phases):
        x = start_x + (card_width + gap) * i
        y = Inches(1.8)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.color.rgb = GCP_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y + Inches(1.0), card_width, Inches(2.6)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = GCP_WHITE
        card.line.color.rgb = RGBColor(0xDA, 0xDC, 0xE0)
        card.line.width = Pt(1)

        # Top accent
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         x, y + Inches(1.0), card_width, Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        # Title
        add_text_box(slide, x + Inches(0.15), y + Inches(1.25),
                     card_width - Inches(0.3), Inches(0.45),
                     title, font_size=16, color=GCP_DARK, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x + Inches(0.15), y + Inches(1.75),
                     card_width - Inches(0.3), Inches(1.5),
                     desc, font_size=12, color=GCP_MEDIUM_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Connecting arrows
    for i in range(4):
        x = start_x + (card_width + gap) * i + card_width
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x - Inches(0.05), Inches(2.0), Inches(0.25), Inches(0.25)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xBD, 0xC1, 0xC6)
        arrow.line.fill.background()


def create_section_divider(prs, section_num, title, subtitle, color=GCP_SECTION_BG):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, color)

    # Large section number
    add_text_box(slide, Inches(1.2), Inches(1.5), Inches(3), Inches(1.5),
                 section_num, font_size=72, color=RGBColor(0xFF, 0xFF, 0xFF),
                 bold=True)

    # Section title
    add_text_box(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(1.2),
                 title, font_size=40, color=GCP_WHITE, bold=True)

    # Subtitle
    add_text_box(slide, Inches(1.2), Inches(4.5), Inches(9), Inches(1.0),
                 subtitle, font_size=18, color=RGBColor(0xAE, 0xCB, 0xFA))

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2), Inches(4.2), Inches(2), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GCP_YELLOW
    line.line.fill.background()

    # Bottom color bar
    colors = [GCP_BLUE, GCP_RED, GCP_YELLOW, GCP_GREEN]
    for i, c in enumerate(colors):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Emu(int(SLIDE_WIDTH / 4 * i)),
                                   SLIDE_HEIGHT - Inches(0.06),
                                   SLIDE_WIDTH // 4, Inches(0.06))
        s.fill.solid()
        s.fill.fore_color.rgb = c
        s.line.fill.background()


def create_content_slide(prs, title, bullets, subtitle="", accent_color=GCP_PRIMARY):
    """Create a standard content slide with title and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, GCP_WHITE)
    add_header_bar(slide, accent_color)
    add_gcp_bar(slide, prs)

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 title, font_size=28, color=GCP_DARK, bold=True)

    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=GCP_MEDIUM_GRAY)
        bullet_y = Inches(1.6)
    else:
        bullet_y = Inches(1.4)

    add_bullet_list(slide, Inches(1.0), bullet_y, Inches(11), Inches(4.5),
                    bullets, font_size=15, spacing=Pt(10))
    return slide


def create_two_column_slide(prs, title, left_title, left_items, right_title, right_items,
                            accent_color=GCP_PRIMARY):
    """Create a two-column content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, GCP_WHITE)
    add_header_bar(slide, accent_color)
    add_gcp_bar(slide, prs)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 title, font_size=28, color=GCP_DARK, bold=True)

    # Left column
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
                 left_title, font_size=18, color=accent_color, bold=True)
    add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(5.3), Inches(4.0),
                    left_items, font_size=13, spacing=Pt(8))

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(6.5), Inches(1.4), Inches(0.02), Inches(5.0))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0xDA, 0xDC, 0xE0)
    div.line.fill.background()

    # Right column
    add_text_box(slide, Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.5),
                 right_title, font_size=18, color=accent_color, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.3), Inches(4.0),
                    right_items, font_size=13, spacing=Pt(8))
    return slide


def create_cards_slide(prs, title, cards, accent_color=GCP_PRIMARY):
    """Create a slide with multiple cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, GCP_LIGHT_GRAY)
    add_header_bar(slide, accent_color)
    add_gcp_bar(slide, prs)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 title, font_size=28, color=GCP_DARK, bold=True)

    cols = min(len(cards), 3)
    card_w = Inches(3.7)
    card_h = Inches(3.5)
    gap = Inches(0.25)
    start_x = Inches(0.5)

    for i, (card_title, card_body, card_color) in enumerate(cards):
        row = i // 3
        col = i % 3
        x = start_x + (card_w + gap) * col
        y = Inches(1.5) + (card_h + gap) * row
        add_card(slide, x, y, card_w, card_h, card_title, card_body, accent_color=card_color)
    return slide


def create_table_slide(prs, title, headers, rows, accent_color=GCP_PRIMARY):
    """Create a slide with a styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, GCP_WHITE)
    add_header_bar(slide, accent_color)
    add_gcp_bar(slide, prs)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 title, font_size=28, color=GCP_DARK, bold=True)

    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(11.5)
    table_height = Inches(min(num_rows * 0.5, 5.0))

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), Inches(1.4),
        table_width, table_height
    )
    table = table_shape.table

    # Style headers
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = GCP_WHITE
            p.font.name = 'Arial'
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent_color

    # Style rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = GCP_DARK
                p.font.name = 'Arial'
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GCP_WHITE
    return slide


def create_key_concepts_slide(prs):
    """Quick Start & Key Concepts slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, GCP_WHITE)
    add_header_bar(slide)
    add_gcp_bar(slide, prs)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Quick Start & Key Concepts", font_size=28, color=GCP_DARK, bold=True)

    # Prerequisites card
    add_card(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(2.5),
             "⚡ Prerequisites",
             "• gcloud CLI\n• terraform\n• kubectl\n• packer\n• jq\n\nAuthenticate & set project,\nthen follow sections in order.",
             accent_color=GCP_BLUE)

    # DWS card
    add_card(slide, Inches(4.6), Inches(1.3), Inches(3.8), Inches(2.5),
             "📊 Dynamic Workload Scheduling",
             "DWS enables Google Cloud to schedule\nworkloads when capacity is available.\n\nCritical for high-demand accelerators\nlike A3/H100.\n\nUp to 53% discount vs on-demand.",
             accent_color=GCP_RED)

    # Hypercompute card
    add_card(slide, Inches(8.7), Inches(1.3), Inches(3.8), Inches(2.5),
             "🚀 GKE AI Hypercompute",
             "Recommended platform for large-scale\nAI training and inference.\n\nIncludes Cluster Toolkit, MIG + DWS\nResize, and Flex Start training.",
             accent_color=GCP_GREEN)

    # Quota card
    add_card(slide, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.5),
             "🔑 Quota & Reservations",
             "Before deploying GPU/TPU workloads, you need:\n• Resource-level quotas (vCPUs, GPUs, machine types)\n"
             "• Reservation count quota (separate from resource quota)\n• Scripts to automate checking & requesting increases",
             accent_color=GCP_YELLOW)

    # Repository structure card
    add_card(slide, Inches(6.6), Inches(4.2), Inches(5.8), Inches(2.5),
             "📁 Repository Structure",
             "01-foundational-tools/     → Agentic coders, accelerator guide\n"
             "02-core-infrastructure/    → Networking, storage, disk images\n"
             "03-deploying-workloads/    → DWS, GKE, Vertex AI, Slurm\n"
             "04-monitoring/             → TPU health, dashboards\n"
             "05-billing-usage/          → Cost controls, CUDs",
             accent_color=GCP_PRIMARY)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ════════════════════════════════════════════════════════
    create_title_slide(prs)

    # ════════════════════════════════════════════════════════
    # SLIDE 2: Agenda
    # ════════════════════════════════════════════════════════
    create_agenda_slide(prs)

    # ════════════════════════════════════════════════════════
    # SLIDE 3: Quick Start & Key Concepts
    # ════════════════════════════════════════════════════════
    create_key_concepts_slide(prs)

    # ════════════════════════════════════════════════════════
    # SECTION 1: FOUNDATIONAL TOOLS
    # ════════════════════════════════════════════════════════
    create_section_divider(prs, "01", "Foundational Tools & Access",
                           "Accelerator selection, agentic coders, quota management")

    # Slide: Accelerator Guide Overview
    create_content_slide(prs,
        "Accelerator Selection & Sizing Guide",
        [
            "🎯  Choosing the right accelerator is the single most impactful decision",
            "💰  GPU hourly rates vary by 10× or more across machine types",
            "⚡  Determines cost, performance, availability, and time-to-result",
            "📐  Use sizing guide BEFORE requesting quota or making reservations",
            "🔄  Covers: GPU catalog, workload matching, memory sizing, checkpointing",
        ],
        subtitle="How to pick the right GPU for your workload",
        accent_color=GCP_BLUE)

    # Slide: GPU Catalog
    create_table_slide(prs,
        "Current-Generation GPU Machine Types",
        ["Machine Type", "GPU Model", "GPUs/VM", "GPU Mem/GPU", "Total GPU Mem", "Network BW"],
        [
            ["A4X Max", "NVIDIA GB300", "4", "279 GB HBM3e", "1,116 GB", "3,600 Gbps"],
            ["A4X", "NVIDIA GB200", "4", "186 GB HBM3e", "744 GB", "2,000 Gbps"],
            ["A4", "NVIDIA B200", "8", "180 GB HBM3e", "1,440 GB", "3,600 Gbps"],
            ["A3 Ultra", "NVIDIA H200", "8", "141 GB HBM3e", "1,128 GB", "3,600 Gbps"],
            ["A3 Mega", "NVIDIA H100", "8", "80 GB HBM3", "640 GB", "1,800 Gbps"],
            ["A3 High (8g)", "NVIDIA H100", "8", "80 GB HBM3", "640 GB", "1,000 Gbps"],
            ["G4 (8 GPU)", "RTX PRO 6000", "8", "96 GB GDDR7", "768 GB", "400 Gbps"],
            ["G2 (8 GPU)", "NVIDIA L4", "8", "24 GB GDDR6", "192 GB", "100 Gbps"],
        ],
        accent_color=GCP_BLUE)

    # Slide: Memory Sizing
    create_table_slide(prs,
        "GPU Memory Sizing — Rules of Thumb",
        ["Scenario", "Approx. GPU Memory Required", "Formula"],
        [
            ["FP16/BF16 Inference", "~2 GB per billion params", "params_B × 2 GB"],
            ["INT8 Inference", "~1 GB per billion params", "params_B × 1 GB"],
            ["INT4 Inference", "~0.5 GB per billion params", "params_B × 0.5 GB"],
            ["BF16 Training (Adam)", "~18-20 GB per billion params", "params_B × 18 GB"],
            ["LoRA Fine-tuning", "~2-4 GB per billion params", "params_B × 2 GB + overhead"],
            ["QLoRA Fine-tuning", "~0.5-1 GB per billion params", "params_B × 0.75 GB + overhead"],
        ],
        accent_color=GCP_BLUE)

    # Slide: Checkpointing
    create_content_slide(prs,
        "Checkpointing — Protect Your Training Investment",
        [
            "⚠️  GPU workloads are ALWAYS at risk of interruption (DWS expiry, Spot preemption, HW failure)",
            "💸  Cost of lost work = hourly_GPU_cost × hours_since_checkpoint × GPU_count",
            "📦  Storage options: Rapid Bucket (best for large scale), GCS (most common), HyperDisk, ParallelStore",
            "🔄  Patterns: Synchronous to GCS (most common), Async with Orbax (JAX/Flax), Rapid Bucket append",
            "⏱️  Frequency: Every 15-30 min (Spot), 30 min-1 hr (DWS), 1-2 hrs (reserved)",
            "⚡  Enable XLA compilation caching for JAX: 52% faster cold starts",
        ],
        subtitle="Never use Local SSD as your only checkpoint destination",
        accent_color=GCP_BLUE)

    # Slide: Agentic Coder Setup
    create_two_column_slide(prs,
        "Agentic Coder Setup",
        "IDE-Based: Cline",
        [
            "Install Cline extension in VS Code",
            "Configure with Vertex AI + Project ID",
            "MCP Servers: Google Dev Knowledge,",
            "  Cloud Logging, GitHub, HuggingFace",
            "Custom rules for coding preferences",
            "Works with Claude Opus 4.6 via Vertex",
        ],
        "CLI-Based: Claude Code & Gemini",
        [
            "Claude Code: npm install globally",
            "  → Configure Vertex AI backend",
            "  → System prompt + permissions",
            "  → Register MCP servers",
            "Gemini CLI: npm install globally",
            "  → Uses gcloud credentials",
            "  → Configure via settings.json",
        ],
        accent_color=GCP_BLUE)

    # ════════════════════════════════════════════════════════
    # SECTION 2: CORE INFRASTRUCTURE
    # ════════════════════════════════════════════════════════
    create_section_divider(prs, "02", "Core Infrastructure Setup",
                           "Networking, disk images, storage, and data pipelines",
                           color=RGBColor(0xC5, 0x22, 0x1F))

    # Slide: Storage for AI Workloads
    create_cards_slide(prs,
        "Storage for AI Workloads",
        [
            ("GCSFuse", "Mount GCS as filesystem via CSI driver.\n7× faster pod startup with tuning.\n41% faster on A3 vs A2.\nStandard pattern in gpu-recipes.", GCP_BLUE),
            ("Rapid Cache", "Managed SSD zonal read cache.\n2.5 TB/s throughput.\nUp to 96% latency reduction.\nZero app changes required.", GCP_RED),
            ("Rapid Bucket", "Zonal object storage with append.\nSub-ms latency, 15 TB/s, 20M QPS.\nBest for model checkpointing.\nCompatible with GCSFuse v3.7.2+.", GCP_GREEN),
        ],
        accent_color=GCP_RED)

    # Slide: Data Pipeline
    create_content_slide(prs,
        "Data Pipeline & Preparation for AI Workloads",
        [
            "📊  BigQuery operates UPSTREAM of the accelerator storage layer",
            "🔄  Common patterns: SQL-to-Training, Dataflow ETL, BigQuery ML, BigFrames notebook",
            "📤  Export formats: Parquet (general), TFRecord (TF), CSV/JSON (small), Avro (schema)",
            "🐼  BigQuery DataFrames: pandas-like API at warehouse scale (750+ APIs)",
            "🔗  Handoff point is GCS: BigQuery prepares → GCS stores → GCSFuse/Rapid delivers to GPU",
            "🏗️  Also supports: Vertex AI Feature Store, Vertex AI Pipelines, Dataflow templates",
        ],
        subtitle="Where BigQuery and data preparation services fit in the AI/ML lifecycle",
        accent_color=GCP_RED)

    # Slide: Disk Images
    create_table_slide(prs,
        "VM Disk Images — 5 Methods for Reusable Configurations",
        ["Method", "Boot Time", "Cost", "Best For"],
        [
            ["Public Images", "Fast", "Free", "Starting point, Deep Learning VMs"],
            ["Custom Images", "Fast", "$0.050/GB/mo", "Golden images, org-wide standards"],
            ["Snapshots (Standard)", "Fast", "$0.026/GB/mo", "Backup, cloning, DR"],
            ["Archive Snapshots", "Slow", "$0.0026/GB/mo", "Compliance, long-term retention"],
            ["Existing Disks", "Instant", "PD pricing", "Quick cloning, dev/test"],
        ],
        accent_color=GCP_RED)

    # Slide: Packer
    create_content_slide(prs,
        "Building VM Images with Packer",
        [
            "🏗️  Infrastructure as Code: image definitions are version-controlled .pkr.hcl files",
            "🔄  Reproducible: every build produces identical image from same template",
            "⚡  One command: packer build template.pkr.hcl (vs. manual multi-step process)",
            "🔌  Provisioners: Shell, Ansible, Chef, File upload — any combination",
            "☁️  CI/CD: Cloud Build integration for automated image pipelines",
            "📦  Examples: Basic Ubuntu, AI/ML GPU (CUDA + PyTorch), HPC (MPI + libfabric)",
            "🔒  Security: Use IAP (no public IP), dedicated service accounts, image scanning",
        ],
        subtitle="Automated, reproducible golden image building on Google Cloud",
        accent_color=GCP_RED)

    # Slide: Zero Trust IAP
    create_two_column_slide(prs,
        "Zero Trust VM Access — No Public IPs + IAP",
        "Benefits",
        [
            "Reduced attack surface — no port scanning",
            "Defense in depth — second layer beyond firewall",
            "No IP-based trust — identity-based access",
            "Simplified firewall management",
            "Compliance: PCI-DSS, HIPAA, SOC 2",
            "Controlled egress via Cloud NAT",
        ],
        "Zero Trust Principles",
        [
            "Verify identity → Google Identity auth",
            "Least privilege → per-user, per-VM IAM",
            "No implicit trust → no VPN, no public IP",
            "Device trust → Access Context Manager",
            "Continuous audit → Cloud Audit Logs",
            "Short-lived creds → OS Login certificates",
        ],
        accent_color=GCP_RED)

    # ════════════════════════════════════════════════════════
    # SECTION 3: DEPLOYING WORKLOADS
    # ════════════════════════════════════════════════════════
    create_section_divider(prs, "03", "Deploying Workloads & Scheduling",
                           "DWS, GKE, Vertex AI, Cluster Toolkit, XPK, Cluster Director",
                           color=RGBColor(0xE3, 0x74, 0x00))

    # Slide: Deployment Overview
    create_table_slide(prs,
        "Deployment Methods at a Glance",
        ["Method", "Platform", "Best For", "DWS Support"],
        [
            ["GKE (Autopilot/Std)", "GKE", "Direct GPU deployment", "Flex-start, Queued"],
            ["Cluster Toolkit", "GKE (Terraform)", "Production clusters, IaC", "Flex-start, Reservation"],
            ["XPK", "GKE (Python CLI)", "Rapid PoC, testing", "Flex-start"],
            ["Cluster Director", "Managed Slurm", "Slurm-native teams, HPC", "Flex-start, Calendar"],
            ["Vertex AI", "Vertex AI", "Serverless training, zero mgmt", "FLEX_START"],
            ["Colab Enterprise", "Vertex AI", "Interactive GPU notebooks", "Via Reservations"],
            ["Compute Engine", "MIGs", "Raw VM control, batch training", "Flex-start"],
            ["Calendar Mode", "Future Reservations", "Guaranteed start time", "Calendar Mode"],
        ],
        accent_color=GCP_YELLOW)

    # Slide: DWS Concepts
    create_two_column_slide(prs,
        "Dynamic Workload Scheduler (DWS)",
        "Flex-Start",
        [
            "Queued until capacity available",
            "Up to 53% discount vs on-demand",
            "Max 7-day run duration",
            "All VMs provisioned at once",
            "Dense allocation (close placement)",
            "Not preemptible once running",
            "Cancel anytime while queued",
        ],
        "Calendar Mode",
        [
            "Guaranteed start time",
            "Must submit 87h+ in advance (GPUs)",
            "Cannot cancel after submission",
            "Committed cost obligation",
            "Dense allocation for GPUs",
            "Up to 53% discount",
            "Best for planned, deadline-critical work",
        ],
        accent_color=GCP_YELLOW)

    # Slide: GKE
    create_two_column_slide(prs,
        "GKE — Autopilot vs Standard Mode",
        "Autopilot (Recommended)",
        [
            "Google manages nodes automatically",
            "Pay per pod resource usage",
            "GPU via nodeSelector (simple!)",
            "Best security defaults enforced",
            "No node pool management needed",
            "DWS via cloud.google.com/gke-flex-start",
            "Limitation: No GPUDirect RDMA",
        ],
        "Standard (Full Control)",
        [
            "You create & manage node pools",
            "Pay for entire VM, used or not",
            "Full control over machine types",
            "GPUDirect RDMA supported",
            "Privileged containers allowed",
            "Custom node configurations",
            "DWS via node pool flags",
        ],
        accent_color=GCP_YELLOW)

    # Slide: Cluster Toolkit
    create_content_slide(prs,
        "Cluster Toolkit — Production-Ready GKE Clusters",
        [
            "🏗️  One-command deployment: gcluster deploy provisions complete stack",
            "📋  Terraform blueprints encode Google's recommended GPU networking settings",
            "🔌  GPUDirect RDMA multi-NIC automatically configured",
            "⚡  Supports reservation-bound, DWS flex-start, and Spot provisioning",
            "📊  Auto-installs Kueue for DWS queued provisioning",
            "🔄  Reservation + DWS fallback pattern via Kueue multi-flavor ClusterQueue",
            "📦  Blueprints: A4X Max, A4X, A4, A3 Ultra, A3 Mega, A3 High",
            "💡  Recommended for production deployments (vs XPK for PoCs)",
        ],
        subtitle="AI-optimized GKE clusters with Terraform blueprints",
        accent_color=GCP_YELLOW)

    # Slide: XPK
    create_content_slide(prs,
        "XPK — Quick GPU & TPU Clusters",
        [
            "🚀  Python CLI that simplifies cluster creation and workload execution",
            "⚡  Zero Kubernetes knowledge required — XPK abstracts all complexity",
            "🔧  One command: xpk cluster create --flex enables DWS flex-start",
            "📊  Auto-installs Kueue, JobSet, CoreDNS, and Crane",
            "🎯  Supports TPU (v4 through Ironwood) and GPU (A100 through A4X)",
            "📊  Built-in priority levels: very-low → very-high with preemption",
            "💡  Best for PoC, testing, experimentation — graduate to Cluster Toolkit for prod",
            "🔄  xpk cluster adapt can configure existing clusters for XPK workloads",
        ],
        subtitle="Accelerated Processing Kit for rapid AI/ML experimentation",
        accent_color=GCP_YELLOW)

    # Slide: Cluster Director
    create_content_slide(prs,
        "Cluster Director — Managed Slurm Clusters",
        [
            "🏢  Fully managed Slurm environment with HA controller, login nodes, GPU compute",
            "🖥️  Console-first: Create clusters via Google Cloud Console step-by-step UI",
            "💰  Supports reservations, DWS Flex-start, Calendar Mode, Spot, On-demand",
            "🔧  Pre-configured: NVIDIA drivers, CUDA, Slurm, NCCL, GPUDirect RDMA pre-installed",
            "🔍  Automatic GPU health checks via Slurm prolog scripts",
            "📊  Topology-aware scheduling for optimal workload placement",
            "🐳  Container support via NVIDIA enroot + pyxis (no custom image needed)",
            "⚡  Minutes to deploy login node + DWS wait for GPU VMs",
        ],
        subtitle="For Slurm-native teams and HPC workloads",
        accent_color=GCP_YELLOW)

    # Slide: Vertex AI & Colab Enterprise
    create_two_column_slide(prs,
        "Vertex AI & Colab Enterprise",
        "Vertex AI (FLEX_START)",
        [
            "Serverless — zero infra management",
            "Fastest time to first job (minutes)",
            "Set strategy to FLEX_START",
            "Supports L4, A100, H100, H200, B200",
            "Max 7-day job timeout",
            "Python SDK, gcloud, or REST API",
            "Up to 53% discount via DWS",
        ],
        "Colab Enterprise",
        [
            "Interactive GPU notebooks",
            "Managed Vertex AI environment",
            "GPUs: V100, T4, A100, L4",
            "Enterprise: IAM, VPC, CMEK",
            "Scheduled notebook execution",
            "Develop → Deploy pattern",
            "18-hour auto-deletion limit",
        ],
        accent_color=GCP_YELLOW)

    # Slide: Compute Engine & Calendar Mode
    create_two_column_slide(prs,
        "Compute Engine & Calendar Mode",
        "MIG Resize Requests (Flex-start)",
        [
            "Raw VM control — no K8s overhead",
            "Create Instance Template → MIG → Resize",
            "All VMs provisioned simultaneously",
            "Max 7-day run duration",
            "Up to 53% discount",
            "Best for batch training, custom images",
            "Teams without K8s expertise",
        ],
        "Calendar Mode (Future Reservations)",
        [
            "Guaranteed start time for GPUs",
            "87-hour minimum lead time",
            "Cannot cancel after submission",
            "Committed cost obligation",
            "Dense allocation for training",
            "Best for planned runs with deadlines",
            "Up to 80 VMs for up to 90 days",
        ],
        accent_color=GCP_YELLOW)

    # ════════════════════════════════════════════════════════
    # SECTION 4: MONITORING & OBSERVABILITY
    # ════════════════════════════════════════════════════════
    create_section_divider(prs, "04", "Monitoring & Observability",
                           "TPU/GPU health, performance profiling, proactive alerting",
                           color=RGBColor(0x0D, 0x65, 0x2D))

    # Slide: Observability Stack
    create_cards_slide(prs,
        "AI Workload Observability Stack",
        [
            ("Layer 1: GKE AI/ML UI", "Single-pane health view:\n• JobSet monitoring dashboard\n• TPU node pool status\n• Accelerator metrics & logs\n• Scheduling & runtime goodput", GCP_BLUE),
            ("Layer 2: Cloud Monitoring", "Metrics & proactive alerts:\n• TPU duty cycle & memory\n• JobSet goodput & TBI/TTR\n• Node health & status\n• PromQL-based alert policies", GCP_RED),
            ("Layer 3: XProf Profiling", "Deep performance analysis:\n• Trace Viewer (operation timeline)\n• HLO Op Profile (time breakdown)\n• Memory Viewer (allocation)\n• Roofline Analysis (bottlenecks)", GCP_GREEN),
        ],
        accent_color=GCP_GREEN)

    # Slide: Key Metrics & Alerts
    create_table_slide(prs,
        "Recommended Proactive Alert Policies",
        ["Alert", "Condition", "Threshold", "Severity"],
        [
            ["TPU Node Not Ready", "condition='Ready', status='False'", "5 min", "Critical"],
            ["Node Pool Error", "status='Error'", "5 min", "Critical"],
            ["Low TPU Duty Cycle", "avg duty cycle < 10%", "15 min", "Warning"],
            ["High HBM Usage", "memory_used/total > 95%", "10 min", "Warning"],
            ["Low Scheduling Goodput", "scheduling_goodput < 0.5", "30 min", "Warning"],
            ["Frequent Interruptions", "interruptions > 3/hour", "1 hour", "Critical"],
        ],
        accent_color=GCP_GREEN)

    # Slide: XProf & TPU Monitoring
    create_two_column_slide(prs,
        "XProf & TPU Monitoring Library",
        "XProf — Performance Profiling",
        [
            "< 1% overhead — profile in production",
            "Trace Viewer: operation timeline on HW",
            "HLO Op Profile: time breakdown",
            "Memory Viewer: allocation visualization",
            "Roofline Analysis: compute vs memory",
            "cloud-diagnostics-xprof for hosting",
            "Supports JAX, PyTorch XLA, TensorFlow",
        ],
        "TPU Monitoring Library (LibTPU)",
        [
            "Hardware-level telemetry from TPU",
            "Metrics: duty_cycle, tensorcore_util",
            "HBM capacity total & usage",
            "Network latency & throughput",
            "HLO execution timing & queue size",
            "TPU-Z for hang/deadlock diagnostics",
            "pip install libtpu or jax[tpu]",
        ],
        accent_color=GCP_GREEN)

    # ════════════════════════════════════════════════════════
    # SECTION 5: BILLING & USAGE
    # ════════════════════════════════════════════════════════
    create_section_divider(prs, "05", "Billing & Usage",
                           "Cost monitoring, budgets, quotas, and optimization",
                           color=GCP_PRIMARY)

    # Slide: Cost Management Stack
    create_cards_slide(prs,
        "Cost Management Stack for AI Workloads",
        [
            ("Visibility", "Cloud Billing Reports, FinOps Hub,\nAnomaly Detection, Gemini Assist.\n\nKnow what you're spending\nacross all projects.", GCP_BLUE),
            ("Guardrails", "Budgets & alerts per project,\nGPU/TPU quotas, programmatic\nauto-disable via Pub/Sub.\n\nPrevent runaway costs.", GCP_RED),
            ("Optimization", "CUDs (up to 57%), DWS (53%),\nSpot (91%), right-sizing,\nidle cleanup, storage tiering.\n\nReduce unit costs.", GCP_GREEN),
        ],
        accent_color=GCP_PRIMARY)

    # Slide: Cost Optimization Strategies
    create_table_slide(prs,
        "GPU/TPU Cost Optimization Strategies",
        ["Strategy", "Savings", "Trade-off", "When to Use"],
        [
            ["Spot / Preemptible VMs", "Up to 91%", "Can be preempted anytime", "Fault-tolerant + checkpointing"],
            ["DWS Flex-Start", "Up to 53%", "Queued, 7-day max", "Time-flexible training jobs"],
            ["CUDs (3-year)", "Up to 57%", "Locked commitment", "Steady-state production"],
            ["CUDs (1-year)", "Up to 46%", "Locked commitment", "Workloads with 1yr+ visibility"],
            ["Right-sizing VMs", "10-40%", "Requires profiling", "Overprovisioned VMs"],
            ["Idle resource cleanup", "100% of idle", "Requires process", "Forgotten clusters/disks"],
        ],
        accent_color=GCP_PRIMARY)

    # Slide: Per-Project Checklist
    create_content_slide(prs,
        "Per-Project Cost Monitoring Checklist",
        [
            "🔴  Create a per-project budget with alerts at 50%, 75%, 90%, 100%",
            "🔴  Enable Billing Export to BigQuery (standard + detailed + pricing)",
            "🔴  Apply labels to ALL resources (team, env, workload, model, cost-center)",
            "🟡  Review and set GPU/TPU quotas — lower quotas = natural spending cap",
            "🟡  Connect Pub/Sub to budget for programmatic alerts",
            "🟡  Enable Anomaly Detection email alerts",
            "🟡  Review FinOps Hub recommendations monthly",
            "🟢  Set up GCS lifecycle policies for checkpoints & profiles",
            "🟢  Evaluate CUD opportunities via FinOps Hub",
        ],
        subtitle="Every project running AI workloads should implement these controls",
        accent_color=GCP_PRIMARY)

    # ════════════════════════════════════════════════════════
    # SUMMARY / THANK YOU
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, GCP_DARK_BLUE)

    # Decorative shapes
    for (x, y, s), col in [
        ((Inches(9.5), Inches(1.0), Inches(2.0)), GCP_BLUE),
        ((Inches(10.8), Inches(3.0), Inches(1.2)), GCP_RED),
        ((Inches(10.2), Inches(4.5), Inches(0.8)), GCP_YELLOW),
        ((Inches(11.5), Inches(2.0), Inches(0.6)), GCP_GREEN),
    ]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, s, s)
        circle.fill.solid()
        circle.fill.fore_color.rgb = col
        circle.line.fill.background()
        circle.fill.fore_color.brightness = 0.15

    add_text_box(slide, Inches(1.2), Inches(1.0), Inches(5), Inches(0.6),
                 "☁  Google Cloud", font_size=22, color=GCP_WHITE, bold=True)

    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(9), Inches(1.0),
                 "Thank You", font_size=48, color=GCP_WHITE, bold=True)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2), Inches(3.5), Inches(3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GCP_YELLOW
    line.line.fill.background()

    add_text_box(slide, Inches(1.2), Inches(3.9), Inches(8), Inches(0.6),
                 "Key Resources", font_size=22, color=RGBColor(0xAE, 0xCB, 0xFA), bold=True)

    resources = [
        "📚  AI Hypercomputer Docs: cloud.google.com/ai-hypercomputer/docs",
        "🔧  Cluster Toolkit: github.com/GoogleCloudPlatform/cluster-toolkit",
        "⚡  XPK: github.com/AI-Hypercomputer/xpk",
        "📊  GPU Recipes: github.com/AI-Hypercomputer/gpu-recipes",
        "💰  DWS Pricing: cloud.google.com/products/dws/pricing",
    ]
    add_bullet_list(slide, Inches(1.2), Inches(4.5), Inches(9), Inches(2.5),
                    resources, font_size=15, color=RGBColor(0xBD, 0xC1, 0xC6),
                    spacing=Pt(12))

    # Bottom bar
    colors = [GCP_BLUE, GCP_RED, GCP_YELLOW, GCP_GREEN]
    for i, c in enumerate(colors):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Emu(int(SLIDE_WIDTH / 4 * i)),
                                   SLIDE_HEIGHT - Inches(0.06),
                                   SLIDE_WIDTH // 4, Inches(0.06))
        s.fill.solid()
        s.fill.fore_color.rgb = c
        s.line.fill.background()

    # ── Save ──
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "AI_Infrastructure_Onboarding_Presentation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
